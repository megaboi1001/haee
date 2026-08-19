#!/usr/bin/env python3
"""Generate the HMSG WEC Racing PRD deliverables with a consistent visual design.

Design system:
  Hyundai-ish corporate palette: navy #002C5F, blue accent #00AAD2,
  amber #F5A623, light backgrounds #F2F7FB / #E8F4FA, gray text #4A4A4A.
"""

from __future__ import annotations

import math
from datetime import date
from pathlib import Path

BASE = Path(__file__).resolve().parent

# ---------------------------------------------------------------------------
# Palette
# ---------------------------------------------------------------------------
NAVY = "002C5F"
BLUE = "007FA8"
SKY = "00AAD2"
AMBER = "F5A623"
GREEN = "1E8449"
RED = "C0392B"
GRAY = "595959"
LIGHT = "F2F7FB"
SKY_FILL = "E8F4FA"
GREEN_FILL = "EAF6EA"
AMBER_FILL = "FDF3E7"
WHITE = "FFFFFF"
BORDER = "BFCDDA"

PROJECT = "HMSG WEC Racing"
REGION = "eu-central-1 (Frankfurt)"
STACK = "hmsg-rac-prd"
TODAY = date.today().strftime("%d %B %Y")

# ===========================================================================
# 1. POWERPOINT  —  AWS Architecture Design Document
# ===========================================================================
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn


def C(hexstr: str) -> RGBColor:
    return RGBColor.from_string(hexstr)


def _set_font(run, size=14, bold=False, color=C(GRAY), italic=False):
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = C(color)
    f.name = "Calibri"


def txt(
    slide,
    x,
    y,
    w,
    h,
    text,
    size=14,
    bold=False,
    color=GRAY,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    wrap=True,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = valign
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        _set_font(r, size=size, bold=bold, color=color)
    return box


def bullets(slide, x, y, w, h, items, size=13.5, color=GRAY, gap=4, marker="\u2022"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    first = True
    for item in items:
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap)
        p.level = level
        if level == 0:
            pPr = p._p.get_or_add_pPr()
            pPr.set("marL", str(182880))
            pPr.set("indent", str(-182880))
            buFont = pPr.makeelement(qn("a:buFont"), {"typeface": "Arial"})
            pPr.append(buFont)
            buChar = pPr.makeelement(qn("a:buChar"), {"char": marker})
            pPr.append(buChar)
        r = p.add_run()
        r.text = "  " + text if level else text
        _set_font(
            r,
            size=size if level == 0 else size - 1.5,
            color=color,
            bold=(level == 0 and False),
        )
    return box


def shape(slide, kind, x, y, w, h, fill=LIGHT, line=BLUE, line_w=1.0, radius=None):
    sp = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = C(fill)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = C(line)
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if radius is not None:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp


def chip(
    slide,
    x,
    y,
    w,
    h,
    text,
    fill=SKY_FILL,
    color=NAVY,
    size=11.5,
    bold=True,
    line=BLUE,
    align=PP_ALIGN.CENTER,
):
    sp = shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill=fill, line=line, radius=0.5
    )
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Pt(3)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    _set_font(r, size=size, bold=bold, color=color)
    return sp


def arrow(slide, x1, y1, x2, y2, color=NAVY, width=2.0, dash=None, both=False):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = C(color)
    conn.line.width = Pt(width)
    if dash:
        from pptx.enum.dml import MSO_LINE

        conn.line.dash_style = MSO_LINE.DASH
    ln = conn.line._get_or_add_ln()
    for tag in ("a:headEnd", "a:tailEnd"):
        for el in ln.findall(qn(tag)):
            ln.remove(el)
    tail = ln.makeelement(
        qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}
    )
    ln.append(tail)
    if both:
        head = ln.makeelement(
            qn("a:headEnd"), {"type": "triangle", "w": "med", "len": "med"}
        )
        ln.append(head)
    return conn


def title_bar(slide, kicker, title):
    txt(slide, 0.55, 0.28, 12.2, 0.3, kicker.upper(), size=11, bold=True, color=SKY)
    txt(slide, 0.55, 0.55, 12.2, 0.62, title, size=29, bold=True, color=NAVY)
    shape(slide, MSO_SHAPE.RECTANGLE, 0.57, 1.28, 1.7, 0.045, fill=AMBER, line=None)


def footer(slide, n, total):
    txt(
        slide,
        0.55,
        7.08,
        9.0,
        0.3,
        f"{PROJECT} · {STACK} · {REGION}",
        size=9,
        color="9AA7B4",
    )
    txt(
        slide,
        11.3,
        7.08,
        1.5,
        0.3,
        f"{n:02d} / {total:02d}",
        size=9,
        color="9AA7B4",
        align=PP_ALIGN.RIGHT,
    )


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
TOTAL = 8

# ---------------------------------------------------------------- Slide 1
s = prs.slides.add_slide(BLANK)
shape(s, MSO_SHAPE.RECTANGLE, -0.02, -0.02, 13.38, 7.54, fill=NAVY, line=None)
shape(s, MSO_SHAPE.RECTANGLE, 0.0, 4.62, 13.38, 0.06, fill=AMBER, line=None)
txt(
    s,
    0.9,
    1.15,
    11.5,
    0.4,
    "PRD · AWS ARCHITECTURE DESIGN DOCUMENT",
    size=13,
    bold=True,
    color=SKY,
)
txt(s, 0.9, 1.55, 11.6, 1.3, "HMSG WEC Racing", size=52, bold=True, color=WHITE)
txt(
    s,
    0.9,
    2.85,
    11.6,
    1.0,
    "Windows MSSQL host · AWS Client VPN (SAML / Entra ID) · Site-to-Site VPN",
    size=19,
    color="D8E4F0",
)
chip(
    s,
    0.9,
    4.0,
    2.6,
    0.5,
    "eu-central-1 · Frankfurt",
    fill=SKY_FILL,
    color=NAVY,
    size=12,
)
chip(s, 3.65, 4.0, 2.35, 0.5, f"Stack {STACK}", fill=SKY_FILL, color=NAVY, size=12)
chip(
    s,
    6.15,
    4.0,
    3.05,
    0.5,
    "Week 1 deliverable",
    fill=AMBER_FILL,
    color=NAVY,
    size=12,
    line=AMBER,
)
txt(
    s,
    0.9,
    6.7,
    11.5,
    0.35,
    f"Prepared by the HAEE Cloud Team · {TODAY}",
    size=11,
    color="9FB4C8",
)

# ---------------------------------------------------------------- Slide 2
s = prs.slides.add_slide(BLANK)
title_bar(s, "Contents", "Agenda")
agenda = [
    ("01", "Requirements at a glance"),
    ("02", "Target architecture"),
    ("03", "Network & security design"),
    ("04", "VPN design — Client VPN · Site-to-Site"),
    ("05", "Deployment & operations"),
    ("06", "Deliverables & next steps"),
]
y = 1.75
for num, text in agenda:
    chip(s, 0.75, y, 0.95, 0.68, num, fill=NAVY, color=WHITE, size=20, line=None)
    txt(s, 2.0, y + 0.14, 9.6, 0.45, text, size=17, bold=True, color=NAVY)
    shape(s, MSO_SHAPE.RECTANGLE, 2.02, y + 0.62, 9.6, 0.014, fill=BORDER, line=None)
    y += 0.82
footer(s, 2, TOTAL)

# ---------------------------------------------------------------- Slide 3
s = prs.slides.add_slide(BLANK)
title_bar(s, "01 · Requirements", "Requirements at a glance")
shape(
    s,
    MSO_SHAPE.ROUNDED_RECTANGLE,
    0.55,
    1.6,
    6.1,
    3.6,
    fill=WHITE,
    line=BORDER,
    radius=0.045,
)
shape(
    s,
    MSO_SHAPE.ROUNDED_RECTANGLE,
    0.55,
    1.6,
    6.1,
    0.16,
    fill=NAVY,
    line=None,
    radius=0.5,
)
txt(s, 0.85, 1.78, 5.5, 0.35, "Scope of this build", size=15, bold=True, color=NAVY)
bullets(
    s,
    0.9,
    2.3,
    5.5,
    2.7,
    [
        "One Windows Server 2022 EC2 instance — m7i.2xlarge (8 vCPU / 32 GiB)",
        "Root 250 GB + data 4 TB (D: · NTFS · 64K) — SQL Server installed by the client",
        "No public IP and no internet egress on the DB host",
        "Private access only: Client VPN for developers, optional Site-to-Site VPN for on-prem",
    ],
    size=13.5,
)
shape(
    s,
    MSO_SHAPE.ROUNDED_RECTANGLE,
    6.95,
    1.6,
    5.85,
    3.6,
    fill=WHITE,
    line=BORDER,
    radius=0.045,
)
shape(
    s,
    MSO_SHAPE.ROUNDED_RECTANGLE,
    6.95,
    1.6,
    5.85,
    0.16,
    fill=BLUE,
    line=None,
    radius=0.5,
)
txt(
    s,
    7.25,
    1.78,
    5.3,
    0.35,
    "Client-side responsibilities",
    size=15,
    bold=True,
    color=NAVY,
)
bullets(
    s,
    7.3,
    2.3,
    5.3,
    2.7,
    [
        "Microsoft Entra ID Enterprise Application (SAML federation)",
        "Entra ID Federation Metadata XML",
        "CustomerGatewayIp + OnPremCidr (enables the Site-to-Site VPN)",
        "SQL Server installation on the EC2 host",
        "VPN user management after handover",
    ],
    size=13.5,
)
shape(
    s,
    MSO_SHAPE.ROUNDED_RECTANGLE,
    0.55,
    5.5,
    12.25,
    1.35,
    fill=AMBER_FILL,
    line=AMBER,
    radius=0.09,
)
txt(s, 0.9, 5.62, 11.7, 0.32, "Constraints & notes", size=12.5, bold=True, color=NAVY)
bullets(
    s,
    0.9,
    5.96,
    11.7,
    0.85,
    [
        "VPC 172.200.0.0/24 is not RFC1918 private space — confirm with the client, or deploy with 172.20.0.0/24 instead.",
        "Up to 10 developers access via Client VPN · Entra ID MFA is the client tenant's policy.",
    ],
    size=11.5,
    gap=2,
)
footer(s, 3, TOTAL)

# ---------------------------------------------------------------- Slide 4
s = prs.slides.add_slide(BLANK)
title_bar(s, "02 · Architecture", "Target architecture")
# AWS boundary
shape(
    s,
    MSO_SHAPE.ROUNDED_RECTANGLE,
    3.15,
    1.5,
    6.95,
    4.95,
    fill="F7FAFD",
    line=NAVY,
    line_w=1.5,
    radius=0.03,
)
txt(
    s,
    3.35,
    1.56,
    6.6,
    0.3,
    "AWS Cloud · VPC hmsg-rac-prd-vpc-ec1 — 172.200.0.0/24",
    size=11.5,
    bold=True,
    color=NAVY,
)
rows = [
    (2.02, SKY_FILL, BLUE, "PUBLIC · pub-ec1a / pub-ec1c — NAT gateway + Elastic IP"),
    (
        2.92,
        GREEN_FILL,
        GREEN,
        "PRIVATE · pri-ec1a / pri-ec1c — Client VPN target networks",
    ),
    (3.82, AMBER_FILL, AMBER, "DB · db-ec1a / db-ec1c — MSSQL host (m7i.2xlarge)"),
]
for y, fill, line, text in rows:
    shape(
        s,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        3.35,
        y,
        6.55,
        0.68,
        fill=fill,
        line=line,
        radius=0.12,
    )
    txt(s, 3.55, y + 0.17, 6.2, 0.4, text, size=12, bold=True, color=NAVY)
# services
txt(
    s,
    3.35,
    5.02,
    6.6,
    0.28,
    "VPC interface endpoints (Secrets Manager · SSM · SSM Messages · EC2 Messages)",
    size=10.5,
    color=GRAY,
)
for i, name in enumerate(["Secrets Manager", "SSM", "CloudWatch"]):
    chip(s, 3.35 + i * 1.66, 5.32, 1.56, 0.5, name, fill=SKY_FILL, color=NAVY, size=11)
txt(
    s,
    3.35,
    5.98,
    6.6,
    0.3,
    "Client VPN logs → CloudWatch (90 d) · AWS API access via VPC interface endpoints only",
    size=10,
    color="9AA7B4",
)
# on-prem
shape(
    s,
    MSO_SHAPE.ROUNDED_RECTANGLE,
    10.45,
    2.6,
    2.55,
    1.05,
    fill="F2F2F2",
    line=GRAY,
    radius=0.08,
)
txt(s, 10.6, 2.72, 2.3, 0.32, "On-prem racing site", size=13, bold=True, color=NAVY)
txt(
    s, 10.6, 3.05, 2.3, 0.5, "Hyundai firewall\nIPsec.1 · static", size=10.5, color=GRAY
)
chip(s, 10.19, 3.32, 0.55, 0.3, "VGW", fill=NAVY, color=WHITE, size=9, line=None)
arrow(s, 10.45, 3.13, 10.62, 3.45, color=GRAY, width=1.75, dash=True)
txt(
    s,
    10.35,
    3.75,
    2.6,
    0.55,
    "Site-to-Site VPN\nonly when on-prem values are supplied",
    size=10,
    color="9AA7B4",
)
# developers + client vpn
shape(
    s,
    MSO_SHAPE.ROUNDED_RECTANGLE,
    0.5,
    2.1,
    2.2,
    0.9,
    fill="F2F2F2",
    line=GRAY,
    radius=0.08,
)
txt(s, 0.65, 2.22, 1.95, 0.3, "Developers (≤10)", size=13, bold=True, color=NAVY)
txt(s, 0.65, 2.55, 1.95, 0.4, "Entra ID SSO · MFA", size=10.5, color=GRAY)
shape(
    s,
    MSO_SHAPE.ROUNDED_RECTANGLE,
    0.5,
    3.6,
    2.2,
    1.1,
    fill=SKY_FILL,
    line=BLUE,
    radius=0.08,
)
txt(s, 0.65, 3.72, 1.95, 0.3, "AWS Client VPN", size=13, bold=True, color=NAVY)
txt(s, 0.65, 4.05, 1.95, 0.55, "UDP 443 · SAML\nsplit tunnel", size=10.5, color=GRAY)
arrow(s, 1.6, 3.0, 1.6, 3.6, color=NAVY, width=1.75)
arrow(s, 2.7, 4.15, 3.35, 3.28, color=NAVY, width=2.0)
txt(s, 1.75, 3.95, 1.7, 0.45, "UDP 443\nSAML SSO", size=9, color=GRAY)
footer(s, 4, TOTAL)

# ---------------------------------------------------------------- Slide 5
s = prs.slides.add_slide(BLANK)
title_bar(s, "03 · Security", "Network & security design")
cards = [
    (
        0.55,
        1.6,
        "Client VPN SG",
        SKY_FILL,
        BLUE,
        [
            "UDP 443 in — 0.0.0.0/0 (endpoint must be reachable)",
            "egress 3389 / 1433 → VPC CIDR (tunnel traffic)",
        ],
    ),
    (
        3.4,
        1.6,
        "Common SG — RDP",
        GREEN_FILL,
        GREEN,
        [
            "TCP 3389 in — 10.200.0.0/22 (Client VPN)",
            "optional: on-prem CIDR via Site-to-Site VPN",
        ],
    ),
    (
        6.25,
        1.6,
        "MSSQL SG — 1433",
        AMBER_FILL,
        AMBER,
        [
            "TCP 1433 in — 10.200.0.0/22 (Client VPN)",
            "optional: on-prem CIDR via Site-to-Site VPN",
        ],
    ),
    (
        9.1,
        1.6,
        "VPC endpoint SG",
        "ECEFF3",
        GRAY,
        ["TCP 443 in — VPC CIDR only", "Response traffic to endpoint ENIs"],
    ),
]
for x, y, title, fill, line, lines in cards:
    shape(
        s,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x,
        y,
        2.6,
        2.15,
        fill=WHITE,
        line=BORDER,
        radius=0.06,
    )
    shape(
        s,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x,
        y,
        2.6,
        0.55,
        fill=fill,
        line=None,
        radius=0.35,
    )
    txt(s, x + 0.18, y + 0.12, 2.3, 0.35, title, size=13, bold=True, color=NAVY)
    bullets(s, x + 0.2, y + 0.72, 2.25, 1.35, lines, size=10.5, gap=4)
shape(
    s,
    MSO_SHAPE.ROUNDED_RECTANGLE,
    0.55,
    4.1,
    12.25,
    2.2,
    fill=LIGHT,
    line=BORDER,
    radius=0.05,
)
txt(
    s,
    0.9,
    4.25,
    11.6,
    0.32,
    "EC2 hardening & guardrails",
    size=12.5,
    bold=True,
    color=NAVY,
)
bullets(
    s,
    0.9,
    4.62,
    11.6,
    1.55,
    [
        "EC2: IMDSv2 required · EBS volumes encrypted · no public IP · termination protected (DeletionPolicy: Retain)",
        "Windows Firewall: RDP 3389 and MSSQL 1433 pre-opened (SQL Server installed later by the client)",
        "Data volume D: formatted NTFS 64K, labelled SQLData — SQL data / log files go here",
        "AWS API access only via VPC interface endpoints — the DB subnet has no NAT / IGW route (by design)",
    ],
    size=11.5,
    gap=3,
)
footer(s, 5, TOTAL)

# ---------------------------------------------------------------- Slide 6
s = prs.slides.add_slide(BLANK)
title_bar(s, "04 · VPN", "VPN design — Client VPN · Site-to-Site")
shape(
    s,
    MSO_SHAPE.ROUNDED_RECTANGLE,
    0.55,
    1.6,
    6.1,
    3.3,
    fill=WHITE,
    line=BORDER,
    radius=0.045,
)
shape(
    s,
    MSO_SHAPE.ROUNDED_RECTANGLE,
    0.55,
    1.6,
    6.1,
    0.16,
    fill=BLUE,
    line=None,
    radius=0.5,
)
txt(s, 0.85, 1.78, 5.5, 0.35, "AWS Client VPN", size=15, bold=True, color=NAVY)
vpn_rows = [
    ("Client CIDR", "10.200.0.0/22"),
    ("Authentication", "Federated SAML — Entra ID"),
    ("Transport / port", "UDP 443 (no TCP fallback)"),
    ("Split tunnel", "Enabled — only VPC traffic"),
    ("Target networks", "pri-ec1a · pri-ec1c"),
    ("Route / authorization", "Whole VPC /24 · all groups"),
    ("Logging", "CloudWatch — 90 days"),
]
y = 2.2
for k, v in vpn_rows:
    txt(s, 0.9, y, 2.35, 0.3, k, size=11, bold=True, color="7C8B99")
    txt(s, 3.3, y, 3.1, 0.3, v, size=11, color=GRAY)
    y += 0.385
shape(
    s,
    MSO_SHAPE.ROUNDED_RECTANGLE,
    6.95,
    1.6,
    5.85,
    3.3,
    fill=WHITE,
    line=BORDER,
    radius=0.045,
)
shape(
    s,
    MSO_SHAPE.ROUNDED_RECTANGLE,
    6.95,
    1.6,
    5.85,
    0.16,
    fill=AMBER,
    line=None,
    radius=0.5,
)
txt(
    s,
    7.25,
    1.78,
    5.3,
    0.35,
    "Site-to-Site VPN (optional)",
    size=15,
    bold=True,
    color=NAVY,
)
bullets(
    s,
    7.3,
    2.3,
    5.25,
    2.5,
    [
        "Created only when the client supplies CustomerGatewayIp + OnPremCidr (add-only stack update)",
        "Customer gateway ipsec.1 · static routing · BGP ASN 65000",
        "Routes propagated to the private and DB route tables",
        "Tunnel configuration downloaded from the VPC console and shared with the on-prem team",
        "If the race site moves (new public IP), the tunnel is replaced and re-configured",
    ],
    size=12,
    gap=4,
)
shape(
    s,
    MSO_SHAPE.ROUNDED_RECTANGLE,
    0.55,
    5.15,
    12.25,
    1.15,
    fill=SKY_FILL,
    line=BLUE,
    radius=0.08,
)
txt(
    s,
    0.9,
    5.3,
    11.6,
    0.3,
    "SAML values for the client's Entra ID Enterprise Application",
    size=12.5,
    bold=True,
    color=NAVY,
)
chips = [
    ("Identifier", "urn:amazon:webservices:clientvpn"),
    ("Reply URL", "https://self-service.clientvpn.amazonaws.com/api/auth/sso/saml"),
    ("Sign-on URL", "http://127.0.0.1:35001 (loopback)"),
]
x = 0.9
for k, v in chips:
    chip(
        s,
        x,
        5.68,
        (11.6 - 0.6) / 3,
        0.5,
        f"{k}: {v}",
        fill=WHITE,
        color=GRAY,
        size=9.5,
        bold=False,
        line=BORDER,
    )
    x += (11.6 - 0.6) / 3 + 0.3
footer(s, 6, TOTAL)

# ---------------------------------------------------------------- Slide 7
s = prs.slides.add_slide(BLANK)
title_bar(s, "05 · Ops", "Deployment & operations")
cols = [
    (
        0.55,
        NAVY,
        "Deploy",
        [
            "CloudFormation stack " + STACK,
            "CAPABILITY_NAMED_IAM required",
            "NoEcho sysadm password → Secrets Manager",
            "Site-to-Site VPN added later via a second deploy — no downtime",
        ],
    ),
    (
        4.65,
        BLUE,
        "Access",
        [
            "RDP via Client VPN to the private IP (sysadm)",
            "SSM Session Manager as backup path",
            "No public IP — internet reachable from VPC only through NAT (private subnets)",
        ],
    ),
    (
        8.75,
        AMBER,
        "Operate",
        [
            "sysadm password rotation: change on host first, then update the secret",
            "Client VPN logs → CloudWatch (90 d)",
            "AWS API access from the DB subnet only via VPC endpoints",
            "Client VPN user management → HMGRS after handover",
        ],
    ),
]
for x, accent, title, items in cols:
    shape(
        s,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x,
        1.6,
        3.85,
        4.3,
        fill=WHITE,
        line=BORDER,
        radius=0.045,
    )
    shape(
        s,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x,
        1.6,
        3.85,
        0.55,
        fill=accent,
        line=None,
        radius=0.35,
    )
    txt(s, x + 0.25, 1.72, 3.4, 0.35, title, size=16, bold=True, color=WHITE)
    bullets(s, x + 0.35, 2.45, 3.25, 3.3, items, size=12, gap=8)
shape(
    s,
    MSO_SHAPE.ROUNDED_RECTANGLE,
    0.55,
    6.15,
    12.25,
    0.6,
    fill=AMBER_FILL,
    line=AMBER,
    radius=0.25,
)
txt(
    s,
    0.9,
    6.3,
    11.6,
    0.35,
    "All temporary implementation IAM permissions are removed at project completion.",
    size=12,
    bold=True,
    color=NAVY,
    align=PP_ALIGN.CENTER,
)
footer(s, 7, TOTAL)

# ---------------------------------------------------------------- Slide 8
s = prs.slides.add_slide(BLANK)
title_bar(s, "06 · Close-out", "Deliverables & next steps")
shape(
    s,
    MSO_SHAPE.ROUNDED_RECTANGLE,
    0.55,
    1.6,
    6.1,
    3.9,
    fill=WHITE,
    line=BORDER,
    radius=0.045,
)
shape(
    s,
    MSO_SHAPE.ROUNDED_RECTANGLE,
    0.55,
    1.6,
    6.1,
    0.16,
    fill=NAVY,
    line=None,
    radius=0.5,
)
txt(
    s,
    0.85,
    1.78,
    5.5,
    0.35,
    "Deliverables (this folder)",
    size=15,
    bold=True,
    color=NAVY,
)
bullets(
    s,
    0.9,
    2.3,
    5.5,
    3.0,
    [
        "AWS Architecture Design Document — PPTX (Week 1)",
        "AWS Infrastructure Resource Inventory — XLSX",
        "Administration Guide — DOCX",
        "Backup · MSSQL EC2 Recovery · VPN Configuration — DOCX",
        "Project Completion Checklist — XLSX",
    ],
    size=13,
    gap=6,
)
shape(
    s,
    MSO_SHAPE.ROUNDED_RECTANGLE,
    6.95,
    1.6,
    5.85,
    3.9,
    fill=WHITE,
    line=BORDER,
    radius=0.045,
)
shape(
    s,
    MSO_SHAPE.ROUNDED_RECTANGLE,
    6.95,
    1.6,
    5.85,
    0.16,
    fill=GREEN,
    line=None,
    radius=0.5,
)
txt(s, 7.25, 1.78, 5.3, 0.35, "Next steps", size=15, bold=True, color=NAVY)
bullets(
    s,
    7.3,
    2.3,
    5.3,
    3.0,
    [
        "Week 1 — review & approve the architecture deck",
        "Client — create Entra ID app and send Federation Metadata XML",
        "HMSG — issue / import ACM certificate, deploy the stack",
        "HMSG — stage SQL media and hand the host over for SQL installation",
        "Close-out — handover docs, remove temporary access, sign-off",
    ],
    size=13,
    gap=6,
)
txt(
    s,
    0.55,
    6.35,
    12.2,
    0.4,
    "Version-controlled in the project repository (main) · generated " + TODAY,
    size=11,
    color="9AA7B4",
)
footer(s, 8, TOTAL)

prs.save(BASE / "AWS-Architecture-Design-Document.pptx")

# ===========================================================================
# 2. EXCEL  —  AWS Infrastructure Resource Inventory
# ===========================================================================
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter
from openpyxl.worksheet.datavalidation import DataValidation
from openpyxl.formatting.rule import FormulaRule
from openpyxl.worksheet.properties import PageSetupProperties

THIN = Side(style="thin", color=BORDER)
HDR_FILL = PatternFill("solid", fgColor=NAVY)
HDR_FONT = Font(name="Calibri", size=11, bold=True, color=WHITE)
TITLE_FONT = Font(name="Calibri", size=18, bold=True, color=NAVY)
SUB_FONT = Font(name="Calibri", size=11, color=GRAY)
BODY_FONT = Font(name="Calibri", size=11, color="333333")
ZEBRA = PatternFill("solid", fgColor="F2F7FB")
CAT_COLOR = {
    "Network": "1F6FB2",
    "Security": "E08A00",
    "Compute": "2E7D32",
    "VPN": "7B4FA6",
}

# --- shared contents ---
inv_headers = [
    "Category",
    "Logical ID",
    "Resource Type",
    "Physical Name / Key Value",
    "Key Settings",
    "Condition",
    "Notes",
]
inv_rows = [
    [
        "Network",
        "Vpc",
        "AWS::EC2::VPC",
        "hmsg-rac-prd-vpc-ec1",
        "CIDR 172.200.0.0/24 · DNS support + hostnames",
        "-",
        "Base VPC",
    ],
    [
        "Network",
        "InternetGateway",
        "AWS::EC2::InternetGateway",
        "hmsg-rac-prd-igw-ec1",
        "Internet gateway",
        "-",
        "Public subnets",
    ],
    [
        "Network",
        "VpcIgwAttachment",
        "AWS::EC2::VPCGatewayAttachment",
        "VPC ↔ IGW",
        "Attachment",
        "-",
        "",
    ],
    [
        "Network",
        "PubSubnetA",
        "AWS::EC2::Subnet",
        "hmsg-rac-prd-sub-pub-ec1a",
        "172.200.0.0/27 · eu-central-1a · public IP on launch",
        "-",
        "NAT host",
    ],
    [
        "Network",
        "PubSubnetC",
        "AWS::EC2::Subnet",
        "hmsg-rac-prd-sub-pub-ec1c",
        "172.200.0.32/27 · eu-central-1c",
        "-",
        "",
    ],
    [
        "Network",
        "PriSubnetA",
        "AWS::EC2::Subnet",
        "hmsg-rac-prd-sub-pri-ec1a",
        "172.200.0.64/27 · eu-central-1a",
        "-",
        "Client VPN target",
    ],
    [
        "Network",
        "PriSubnetC",
        "AWS::EC2::Subnet",
        "hmsg-rac-prd-sub-pri-ec1c",
        "172.200.0.96/27 · eu-central-1c",
        "-",
        "Client VPN target",
    ],
    [
        "Network",
        "DbSubnetA",
        "AWS::EC2::Subnet",
        "hmsg-rac-prd-sub-db-ec1a",
        "172.200.0.128/27 · eu-central-1a",
        "-",
        "MSSQL host + endpoints",
    ],
    [
        "Network",
        "DbSubnetC",
        "AWS::EC2::Subnet",
        "hmsg-rac-prd-sub-db-ec1c",
        "172.200.0.160/27 · eu-central-1c",
        "-",
        "",
    ],
    [
        "Network",
        "PubRouteTable",
        "AWS::EC2::RouteTable",
        "hmsg-rac-prd-rtb-pub-ec1",
        "0.0.0.0/0 → IGW",
        "-",
        "",
    ],
    [
        "Network",
        "PriRouteTable",
        "AWS::EC2::RouteTable",
        "hmsg-rac-prd-rtb-pri-ec1",
        "0.0.0.0/0 → NAT · VGW propagation",
        "-",
        "",
    ],
    [
        "Network",
        "DbRouteTable",
        "AWS::EC2::RouteTable",
        "hmsg-rac-prd-rtb-db-ec1",
        "local only · VGW propagation · no IGW/NAT",
        "-",
        "",
    ],
    [
        "Network",
        "PubRouteInternet",
        "AWS::EC2::Route",
        "0.0.0.0/0 → IGW",
        "Public route",
        "-",
        "",
    ],
    [
        "Network",
        "PriRouteNat",
        "AWS::EC2::Route",
        "0.0.0.0/0 → NAT",
        "Private outbound",
        "-",
        "",
    ],
    [
        "Network",
        "NatEip",
        "AWS::EC2::EIP",
        "NAT Elastic IP",
        "Allocated for NAT gateway",
        "-",
        "",
    ],
    [
        "Network",
        "NatGateway",
        "AWS::EC2::NatGateway",
        "hmsg-rac-prd-NAT-pub-ec1",
        "Public subnet ec1a",
        "-",
        "",
    ],
    [
        "Network",
        "PubSubnetARouteAssoc",
        "AWS::EC2::SubnetRouteTableAssociation",
        "pub-a → pub rtb",
        "Association",
        "-",
        "",
    ],
    [
        "Network",
        "PubSubnetCRouteAssoc",
        "AWS::EC2::SubnetRouteTableAssociation",
        "pub-c → pub rtb",
        "Association",
        "-",
        "",
    ],
    [
        "Network",
        "PriSubnetARouteAssoc",
        "AWS::EC2::SubnetRouteTableAssociation",
        "pri-a → pri rtb",
        "Association",
        "-",
        "",
    ],
    [
        "Network",
        "PriSubnetCRouteAssoc",
        "AWS::EC2::SubnetRouteTableAssociation",
        "pri-c → pri rtb",
        "Association",
        "-",
        "",
    ],
    [
        "Network",
        "DbSubnetARouteAssoc",
        "AWS::EC2::SubnetRouteTableAssociation",
        "db-a → db rtb",
        "Association",
        "-",
        "",
    ],
    [
        "Network",
        "DbSubnetCRouteAssoc",
        "AWS::EC2::SubnetRouteTableAssociation",
        "db-c → db rtb",
        "Association",
        "-",
        "",
    ],
    [
        "Security",
        "CommonSg",
        "AWS::EC2::SecurityGroup",
        "hmsg-rac-prd-common-ec2-sg-ec1",
        "RDP 3389 ← 10.200.0.0/22 · egress 443 → VPC",
        "-",
        "Windows admin",
    ],
    [
        "Security",
        "MssqlSg",
        "AWS::EC2::SecurityGroup",
        "hmsg-rac-prd-mssql-ec2-sg-ec1",
        "MSSQL 1433 ← 10.200.0.0/22 · egress 443 → VPC",
        "-",
        "SQL access",
    ],
    [
        "Security",
        "CommonSgOnPremRdp",
        "AWS::EC2::SecurityGroupIngress",
        "RDP from on-prem",
        "3389 ← OnPremCidr",
        "CreateSiteToSiteVpn",
        "",
    ],
    [
        "Security",
        "MssqlSgOnPremMssql",
        "AWS::EC2::SecurityGroupIngress",
        "MSSQL from on-prem",
        "1433 ← OnPremCidr",
        "CreateSiteToSiteVpn",
        "",
    ],
    [
        "Security",
        "ClientVpnSg",
        "AWS::EC2::SecurityGroup",
        "hmsg-rac-prd-clientvpn-sg-ec1",
        "UDP 443 ← 0.0.0.0/0 · egress 3389/1433 → VPC",
        "-",
        "Endpoint SG",
    ],
    [
        "Security",
        "VpcEndpointSg",
        "AWS::EC2::SecurityGroup",
        "hmsg-rac-prd-vpce-sg-ec1",
        "TCP 443 ← VPC CIDR",
        "-",
        "Interface endpoints",
    ],
    [
        "Compute",
        "SysadmSecret",
        "AWS::SecretsManager::Secret",
        "hmsg-rac-prd-sysadm-password",
        "sysadm password (NoEcho) · retained",
        "-",
        "",
    ],
    [
        "Compute",
        "Ec2Role",
        "AWS::IAM::Role",
        "EC2 instance role",
        "AmazonSSMManagedInstanceCore · GetSecretValue (1 secret)",
        "-",
        "",
    ],
    [
        "Compute",
        "Ec2InstanceProfile",
        "AWS::IAM::InstanceProfile",
        "EC2 instance profile",
        "Attaches Ec2Role",
        "-",
        "",
    ],
    [
        "Compute",
        "MssqlInstance",
        "AWS::EC2::Instance",
        "hmsg-rac-prd-mssql01-ec2-ec1a",
        "m7i.2xlarge · Win 2022 EN · 250 GB root + 4 TB D: · IMDSv2",
        "-",
        "EB S-optimised · retained",
    ],
    [
        "Compute",
        "SecretsManagerEndpoint",
        "AWS::EC2::VPCEndpoint",
        "com.amazonaws.eu-central-1.secretsmanager",
        "Interface · DbSubnetA · private DNS",
        "-",
        "",
    ],
    [
        "Compute",
        "SsmEndpoint",
        "AWS::EC2::VPCEndpoint",
        "com.amazonaws.eu-central-1.ssm",
        "Interface · DbSubnetA · private DNS",
        "-",
        "",
    ],
    [
        "Compute",
        "SsmMessagesEndpoint",
        "AWS::EC2::VPCEndpoint",
        "com.amazonaws.eu-central-1.ssmmessages",
        "Interface · DbSubnetA · private DNS",
        "-",
        "",
    ],
    [
        "Compute",
        "Ec2MessagesEndpoint",
        "AWS::EC2::VPCEndpoint",
        "com.amazonaws.eu-central-1.ec2messages",
        "Interface · DbSubnetA · private DNS",
        "-",
        "",
    ],
    [
        "VPN",
        "SamlProvider",
        "AWS::IAM::SAMLProvider",
        "hmsg-rac-prd-clientvpn-saml",
        "Entra ID Federation Metadata XML",
        "-",
        "",
    ],
    [
        "VPN",
        "AcmCertificate",
        "AWS::CertificateManager::Certificate",
        "vpn.hmg-racing.com",
        "Public DNS validation (hosted zone)",
        "CreateAcmCertHostedZone",
        "",
    ],
    [
        "VPN",
        "AcmCertificateManual",
        "AWS::CertificateManager::Certificate",
        "vpn.hmg-racing.com",
        "Public DNS validation (manual CNAME)",
        "CreateAcmCertManual",
        "",
    ],
    [
        "VPN",
        "ClientVpnLogGroup",
        "AWS::Logs::LogGroup",
        "hmsg-rac-prd-clientvpn-logs",
        "CloudWatch · 90 day retention",
        "-",
        "",
    ],
    [
        "VPN",
        "ClientVpnEndpoint",
        "AWS::EC2::ClientVpnEndpoint",
        "hmsg-rac-prd-clientvpn-ec1",
        "CIDR 10.200.0.0/22 · SAML federation · UDP 443 · split tunnel",
        "-",
        "",
    ],
    [
        "VPN",
        "ClientVpnAssocPriA",
        "AWS::EC2::ClientVpnTargetNetworkAssociation",
        "pri-a association",
        "Target network PriSubnetA",
        "-",
        "",
    ],
    [
        "VPN",
        "ClientVpnAssocPriC",
        "AWS::EC2::ClientVpnTargetNetworkAssociation",
        "pri-c association",
        "Target network PriSubnetC",
        "-",
        "",
    ],
    [
        "VPN",
        "ClientVpnRouteVpc",
        "AWS::EC2::ClientVpnRoute",
        "route → VPC /24",
        "Via PriSubnetA association",
        "-",
        "",
    ],
    [
        "VPN",
        "ClientVpnAuthRule",
        "AWS::EC2::ClientVpnAuthorizationRule",
        "authorize all groups",
        "TargetNetworkCidr = VPC /24",
        "-",
        "",
    ],
    [
        "VPN",
        "CustomerGateway",
        "AWS::EC2::CustomerGateway",
        "hmsg-rac-prd-cgw-ec1",
        "ipsec.1 · BGP ASN 65000 · public IP from client",
        "CreateSiteToSiteVpn",
        "",
    ],
    [
        "VPN",
        "VpnGateway",
        "AWS::EC2::VPNGateway",
        "hmsg-rac-prd-vgw-ec1",
        "ipsec.1 · Amazon side ASN 64512",
        "CreateSiteToSiteVpn",
        "",
    ],
    [
        "VPN",
        "VpcVgwAttachment",
        "AWS::EC2::VPCGatewayAttachment",
        "VPC ↔ VGW",
        "Attachment",
        "CreateSiteToSiteVpn",
        "",
    ],
    [
        "VPN",
        "VpnConnection",
        "AWS::EC2::VPNConnection",
        "hmsg-rac-prd-vpn-conn-ec1",
        "static routes only · 2 tunnels",
        "CreateSiteToSiteVpn",
        "",
    ],
    [
        "VPN",
        "VpnConnectionRouteOnPrem",
        "AWS::EC2::VPNConnectionRoute",
        "on-prem route",
        "Destination = OnPremCidr",
        "CreateSiteToSiteVpn",
        "",
    ],
    [
        "VPN",
        "VpnPropagationPri",
        "AWS::EC2::VPNGatewayRoutePropagation",
        "propagate → pri rtb",
        "VGW propagation",
        "CreateSiteToSiteVpn",
        "",
    ],
    [
        "VPN",
        "VpnPropagationDb",
        "AWS::EC2::VPNGatewayRoutePropagation",
        "propagate → db rtb",
        "VGW propagation",
        "CreateSiteToSiteVpn",
        "",
    ],
]

param_rows = [
    ["ProjectName", "String", "hmsg-rac", "Name prefix"],
    ["EnvironmentName", "String", "prd", "Environment suffix (prd only)"],
    [
        "VpcCidr",
        "String",
        "172.200.0.0/24",
        "VPC CIDR — /24; subnets derived via Fn::Cidr",
    ],
    ["ClientVpnCidr", "String", "10.200.0.0/22", "Client VPN client pool"],
    ["OnPremCidr", "String", "(empty)", "Optional — blank skips the Site-to-Site VPN"],
    [
        "CustomerGatewayIp",
        "String",
        "(empty)",
        "Optional — blank skips the Site-to-Site VPN",
    ],
    ["InstanceType", "String", "m7i.2xlarge", "Windows host (8 vCPU / 32 GiB)"],
    [
        "KeyPairName",
        "AWS::EC2::KeyPair::KeyName",
        "hmsg-rac-prd-key-ec2-ec1",
        "Admin password decrypt",
    ],
    [
        "WindowsAmi",
        "AWS::SSM::Parameter::Value<AWS::EC2::Image::Id>",
        "/aws/service/ami-windows-latest/Windows_Server-2022-English-Full-Base",
        "Latest Windows 2022 EN",
    ],
    [
        "AdminPassword",
        "String (NoEcho)",
        "(hidden)",
        "sysadm password → Secrets Manager",
    ],
    ["AdminSecretName", "String", "hmsg-rac-prd-sysadm-password", "Secret name"],
    ["VpnDomainName", "String", "vpn.hmg-racing.com", "Server certificate subject"],
    ["HostedZoneId", "String", "(empty)", "Optional Route 53 zone for ACM validation"],
    ["ServerCertificateArn", "String", "(empty)", "Optional existing ACM certificate"],
    ["SplitTunnel", "String", "true", "true / false"],
]

cond_rows = [
    ["UseExistingCert", "ServerCertificateArn provided?", "Reuse supplied cert"],
    ["UseHostedZone", "HostedZoneId provided?", "Automatic ACM DNS validation"],
    ["NoHostedZone", "HostedZoneId empty", "Manual DNS validation"],
    ["SplitTunnelEnabled", "SplitTunnel = true", "Split tunnel on"],
    [
        "CreateAcmCertHostedZone",
        "UseHostedZone AND NOT UseExistingCert",
        "New cert in hosted zone",
    ],
    [
        "CreateAcmCertManual",
        "NoHostedZone AND NOT UseExistingCert",
        "New cert, manual CNAME",
    ],
    ["HasOnPremCidr", "OnPremCidr provided", ""],
    ["HasCustomerGatewayIp", "CustomerGatewayIp provided", ""],
    [
        "CreateSiteToSiteVpn",
        "HasOnPremCidr AND HasCustomerGatewayIp",
        "Build the S2S VPN",
    ],
]

out_rows = [
    ["VpcId", "VPC ID", ""],
    ["NatEip", "NAT Elastic IP", "Public"],
    ["DbSubnetA / DbSubnetC", "DB subnet IDs", ""],
    ["MssqlInstanceId", "EC2 instance ID", ""],
    ["MssqlPrivateIp", "Private IPv4", "RDP / SQL target"],
    ["SysadmSecretArn", "Secrets Manager ARN", "sysadm password"],
    ["ServerCertificateArn", "ACM certificate ARN", "Client VPN server cert"],
    ["SamlProviderArn", "IAM SAML provider ARN", ""],
    ["ClientVpnEndpointId", "Endpoint ID", ""],
    ["ClientVpnLogGroup", "CloudWatch log group", ""],
    ["CustomerGatewayId", "CGW ID", "Only with S2S VPN"],
    ["VpnConnectionId", "VPN connection ID", "Tunnel config source"],
    ["SiteToSiteVpnDeployed", "true / false", "Indicator"],
]

vpn_rows = [
    ["Client VPN", "Client CIDR", "10.200.0.0/22", ""],
    ["Client VPN", "Authentication", "Federated SAML (Entra ID)", "Client-side app"],
    ["Client VPN", "Transport", "UDP 443", "No TCP fallback"],
    ["Client VPN", "Split tunnel", "Enabled", ""],
    ["Client VPN", "Target networks", "PriSubnetA · PriSubnetC", ""],
    ["Client VPN", "Route", "172.200.0.0/24 via pri-a", "Whole VPC"],
    ["Client VPN", "Authorization", "All groups", "Restrict via Entra group later"],
    [
        "Client VPN",
        "Self-service portal",
        "https://self-service.clientvpn.amazonaws.com",
        "SAML SSO",
    ],
    [
        "SAML values",
        "Identifier",
        "urn:amazon:webservices:clientvpn",
        "For the Entra ID app",
    ],
    [
        "SAML values",
        "Reply URL",
        "https://self-service.clientvpn.amazonaws.com/api/auth/sso/saml",
        "",
    ],
    ["SAML values", "Sign-on URL", "http://127.0.0.1:35001", "Loopback on user device"],
    [
        "SAML metadata",
        "How it is applied",
        "Out-of-band: aws iam update-saml-provider",
        "Real Entra ID XML is ~13 KB — too large for a CFN String param (4096 limit); template carries a placeholder provider",
    ],
    [
        "SAML metadata",
        "Post-deploy command",
        "update-saml-provider --saml-provider-arn <output> --saml-metadata-document file://<xml>",
        "Before anyone can SSO in",
    ],
    ["Site-to-Site", "Customer gateway", "ipsec.1 · ASN 65000", "Conditional"],
    ["Site-to-Site", "Routing", "Static · prop. to pri/db RTs", ""],
    [
        "Site-to-Site",
        "Tunnel config",
        "VPC console → download configuration",
        "Share with on-prem team",
    ],
]

perm_rows = [
    [
        "Client VPN users (devs)",
        "Client-side",
        "Entra ID group the client assigns",
        "Access only through Client VPN — no public reach",
    ],
    [
        "Implementation team (HMSG)",
        "Temporary",
        "AWS console / IAM admin",
        "Removed at project completion",
    ],
    [
        "EC2 host operators",
        "Secure ops",
        "sysadm (Secrets Manager) · SSM Session Manager",
        "RDP via Client VPN only",
    ],
    [
        "On-prem network / firewall",
        "Client infra",
        "Customer gateway",
        "Manages the Site-to-Site VPN side",
    ],
]


def style_sheet(ws, headers, rows, widths, zebra=True, tab=None, wrap_cols=None):
    if tab:
        ws.sheet_properties.tabColor = tab
    ws.append(headers)
    for c in ws[1]:
        c.fill = HDR_FILL
        c.font = HDR_FONT
        c.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        c.border = Border(left=THIN, right=THIN, top=THIN, bottom=THIN)
    ws.row_dimensions[1].height = 30
    for r in rows:
        ws.append(r)
    wrap_cols = wrap_cols or list(range(1, len(headers) + 1))
    for i, row in enumerate(ws.iter_rows(min_row=2), start=2):
        for j, c in enumerate(row, start=1):
            c.font = BODY_FONT
            c.border = Border(left=THIN, right=THIN, top=THIN, bottom=THIN)
            c.alignment = Alignment(
                vertical="top",
                wrap_text=(j in wrap_cols),
                horizontal="center" if j == 1 else "left",
            )
            if zebra and i % 2 == 0:
                c.fill = ZEBRA
        est = 1
        for j, v in enumerate(row, start=1):
            if v is None:
                continue
            cw = widths.get(get_column_letter(j), 20)
            est = max(est, math.ceil(len(str(v)) / max(8, cw)))
        ws.row_dimensions[i].height = min(90, max(17, est * 14.5))
    for col, w in widths.items():
        ws.column_dimensions[col].width = w
    ws.freeze_panes = "A2"
    ws.auto_filter.ref = f"A1:{get_column_letter(len(headers))}{ws.max_row}"
    ws.page_setup.orientation = "landscape"
    ws.sheet_properties.pageSetUpPr = PageSetupProperties(fitToPage=True)
    ws.page_setup.fitToWidth = 1
    ws.page_setup.fitToHeight = 0


wb = Workbook()

# Summary sheet
ws = wb.active
ws.title = "Summary"
ws.sheet_properties.tabColor = NAVY
ws.sheet_view.showGridLines = False
ws.merge_cells("A1:F1")
ws["A1"] = "HMSG WEC Racing — PRD"
ws["A1"].font = Font(name="Calibri", size=22, bold=True, color=NAVY)
ws.merge_cells("A2:F2")
ws["A2"] = "AWS Infrastructure Resource Inventory & VPN Configuration"
ws["A2"].font = Font(name="Calibri", size=13, color=BLUE)
ws.row_dimensions[1].height = 34
ws.row_dimensions[2].height = 20
meta = [
    ("Document", "AWS Infrastructure Resource Inventory"),
    ("Project", "HMSG (Hyundai Motorsport GmbH) — WEC Racing"),
    ("Environment", "Production (prd)"),
    ("Region", "eu-central-1 (Frankfurt)"),
    ("Stack", "hmsg-rac-prd"),
    ("CloudFormation", "hmsg-rac-prd.yaml"),
    ("Generated", TODAY),
    ("Baseline", "main @ def9a32 — deliverables commit"),
]
for i, (k, v) in enumerate(meta, start=4):
    ws.cell(row=i, column=1, value=k).font = Font(
        name="Calibri", size=11, bold=True, color=NAVY
    )
    ws.cell(row=i, column=1).fill = PatternFill("solid", fgColor=SKY_FILL)
    ws.cell(row=i, column=1).border = Border(
        left=THIN, right=THIN, top=THIN, bottom=THIN
    )
    ws.merge_cells(start_row=i, start_column=2, end_row=i, end_column=5)
    c = ws.cell(row=i, column=2, value=v)
    c.font = BODY_FONT
    c.border = Border(left=THIN, right=THIN, top=THIN, bottom=THIN)
    ws.row_dimensions[i].height = 20
ws.merge_cells("A14:F15")
ws["A14"] = (
    "Sheets: Resource Inventory (54 resources) · Parameters (16) · Conditions (9) · "
    "Outputs (14) · VPN Config · Users & Permissions."
)
ws["A14"].font = Font(name="Calibri", size=11, italic=True, color=GRAY)
ws["A14"].alignment = Alignment(wrap_text=True, vertical="top")
ws.row_dimensions[14].height = 30
for col, w in {"A": 18, "B": 22, "C": 22, "D": 22, "E": 22, "F": 22}.items():
    ws.column_dimensions[col].width = w

# Resource Inventory
ws = wb.create_sheet("Resource Inventory")
style_sheet(
    ws,
    inv_headers,
    inv_rows,
    {"A": 12, "B": 26, "C": 34, "D": 32, "E": 52, "F": 20, "G": 24},
    tab=BLUE,
)
for i, row in enumerate(ws.iter_rows(min_row=2, max_row=ws.max_row), start=2):
    cat = row[0].value
    row[0].font = Font(
        name="Calibri", size=11, bold=True, color=CAT_COLOR.get(cat, NAVY)
    )

# Parameters
ws = wb.create_sheet("Parameters")
style_sheet(
    ws,
    ["Parameter", "Type", "Default / Example", "Purpose"],
    param_rows,
    {"A": 24, "B": 40, "C": 46, "D": 46},
    tab=SKY,
)

# Conditions
ws = wb.create_sheet("Conditions")
style_sheet(
    ws,
    ["Condition", "Expression (compact)", "Effect"],
    cond_rows,
    {"A": 28, "B": 62, "C": 46},
    tab="9AA7B4",
)

# Outputs
ws = wb.create_sheet("Outputs")
style_sheet(
    ws,
    ["Output", "What it returns", "Use"],
    out_rows,
    {"A": 28, "B": 40, "C": 40},
    tab=AMBER,
)

# VPN Config
ws = wb.create_sheet("VPN Config")
style_sheet(
    ws,
    ["Area", "Setting", "Value", "Notes"],
    vpn_rows,
    {"A": 16, "B": 24, "C": 52, "D": 34},
    tab=SKY,
)

# Users & Permissions
ws = wb.create_sheet("Users & Permissions")
style_sheet(
    ws,
    ["Role / Owner", "Scope", "Access", "Notes"],
    perm_rows,
    {"A": 26, "B": 16, "C": 44, "D": 40},
    tab=GREEN,
)

wb.save(BASE / "AWS-Infrastructure-Resource-Inventory.xlsx")

# ---------------------------------------------------------------------------
# 3. EXCEL — Project Completion Checklist
# ---------------------------------------------------------------------------
cwb = Workbook()
cs = cwb.active
cs.title = "Checklist"
cs.sheet_properties.tabColor = NAVY
chk_headers = ["Category", "Checklist item", "Owner", "Status", "Evidence / notes"]
chk_rows = [
    [
        "Requirements",
        "Confirm VpcCidr (172.200.0.0/24) with the client — non-RFC1918 range",
        "Client / HMSG",
    ],
    ["Requirements", "Confirm ClientVpnCidr 10.200.0.0/22 has no overlap", "HMSG"],
    [
        "Requirements",
        "Receive CustomerGatewayIp + OnPremCidr (for the optional S2S VPN)",
        "Client",
    ],
    ["Requirements", "Receive Entra ID Federation Metadata XML", "Client"],
    ["Design", "Architecture design document approved (Week 1)", "Client / HMSG"],
    ["Network", "Deploy VPC / subnets / route tables / IGW / NAT", "HMSG"],
    [
        "Security",
        "Security groups allow RDP / MSSQL only from VPN + on-prem ranges",
        "HMSG",
    ],
    ["Security", "sysadm secret exists in Secrets Manager (NoEcho)", "HMSG"],
    ["Compute", "EC2 launched in db-ec1a (m7i.2xlarge · Win 2022)", "HMSG"],
    ["Compute", "D: 4 TB volume formatted NTFS and mounted", "HMSG"],
    ["Compute", "sysadm local account created and tested", "HMSG"],
    ["VPN", "ACM server certificate issued / imported", "HMSG"],
    ["VPN", "IAM SAML provider updated with real metadata XML", "HMSG"],
    ["VPN", "Client VPN endpoint validated; .ovpn distributed to ≤10 users", "HMSG"],
    [
        "VPN",
        "Developer access to the DB host over Client VPN verified",
        "Client / HMSG",
    ],
    ["VPN", "Site-to-Site VPN created and tunnels UP (if required)", "Client / HMSG"],
    ["VPN", "On-prem route propagated into pri/db route tables verified", "HMSG"],
    ["Docs", "Resource inventory delivered", "HMSG"],
    ["Docs", "Administration guide delivered", "HMSG"],
    ["Docs", "Backup / recovery / VPN configuration delivered", "HMSG"],
    ["Docs", "Completion checklist closed out", "HMSG"],
    ["Handover", "Temporary implementation IAM access removed", "HMSG"],
    ["Handover", "VPN user management transferred to HMGRS", "Client / HMGRS"],
    ["Handover", "Final sign-off obtained", "Client / HMSG"],
]
style_sheet(
    cs,
    chk_headers,
    [r + ["Not started", ""] for r in chk_rows],
    {"A": 16, "B": 68, "C": 16, "D": 15, "E": 40},
    zebra=True,
    tab=NAVY,
)
dv = DataValidation(
    type="list", formula1='"Not started,In progress,Blocked,Done"', allow_blank=True
)
dv.error = "Choose: Not started, In progress, Blocked or Done."
dv.errorTitle = "Invalid status"
cs.add_data_validation(dv)
dv.add(f"D2:D{cs.max_row}")
rng = f"D2:D{cs.max_row}"
cs.conditional_formatting.add(
    rng,
    FormulaRule(
        formula=['$D2="Done"'],
        fill=PatternFill("solid", fgColor="C6EFCE"),
        font=Font(color="006100", bold=True),
    ),
)
cs.conditional_formatting.add(
    rng,
    FormulaRule(
        formula=['$D2="In progress"'],
        fill=PatternFill("solid", fgColor="FFEB9C"),
        font=Font(color="9C6500", bold=True),
    ),
)
cs.conditional_formatting.add(
    rng,
    FormulaRule(
        formula=['$D2="Blocked"'],
        fill=PatternFill("solid", fgColor="FFC7CE"),
        font=Font(color="9C0006", bold=True),
    ),
)
cs.conditional_formatting.add(
    rng,
    FormulaRule(
        formula=['$D2="Not started"'],
        fill=PatternFill("solid", fgColor="EDEDED"),
        font=Font(color="595959"),
    ),
)

ls = cwb.create_sheet("Legend")
ls.sheet_properties.tabColor = "9AA7B4"
ls.sheet_view.showGridLines = False
ls.merge_cells("A1:C1")
ls["A1"] = "Status legend"
ls["A1"].font = Font(name="Calibri", size=16, bold=True, color=NAVY)
legend = [
    ("Done", "C6EFCE", "006100"),
    ("In progress", "FFEB9C", "9C6500"),
    ("Blocked", "FFC7CE", "9C0006"),
    ("Not started", "EDEDED", "595959"),
]
for i, (label, fill, txtc) in enumerate(legend, start=3):
    c = ls.cell(row=i, column=1, value=label)
    c.fill = PatternFill("solid", fgColor=fill)
    c.font = Font(name="Calibri", size=12, bold=True, color=txtc)
    c.border = Border(left=THIN, right=THIN, top=THIN, bottom=THIN)
    ls.column_dimensions["A"].width = 20
ls["A9"] = "Owner: HMSG = implementation team · Client = Hyundai / HMGRS"
ls["A9"].font = Font(name="Calibri", size=10, italic=True, color=GRAY)
cwb.save(BASE / "Project-Completion-Checklist.xlsx")

# ===========================================================================
# 4. WORD  —  Administration Guide + Backup / Recovery / VPN Configuration
# ===========================================================================
from docx import Document
from docx.shared import Inches, Pt, RGBColor as WColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml import OxmlElement
from docx.oxml.ns import qn


def shade(cell, hexfill):
    tcPr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:val"), "clear")
    shd.set(qn("w:fill"), hexfill)
    tcPr.append(shd)


def add_bottom_border(paragraph, color=BLUE, size="14"):
    pPr = paragraph._p.get_or_add_pPr()
    pBdr = OxmlElement("w:pBdr")
    bottom = OxmlElement("w:bottom")
    bottom.set(qn("w:val"), "single")
    bottom.set(qn("w:sz"), size)
    bottom.set(qn("w:space"), "4")
    bottom.set(qn("w:color"), color)
    pBdr.append(bottom)
    pPr.append(pBdr)


def style_heading(doc, level, color=NAVY):
    st = doc.styles[f"Heading {level}"]
    st.font.name = "Calibri"
    st.font.color.rgb = WColor(0, 44, 95)
    st.font.size = Pt(15 if level == 1 else (13 if level == 2 else 11.5))
    st.paragraph_format.space_before = Pt(14 if level == 1 else 8)
    st.paragraph_format.space_after = Pt(4)
    st.paragraph_format.line_spacing = 1.05


def doc_table(doc, headers, rows, widths=None, zebra=True):
    t = doc.add_table(rows=1, cols=len(headers))
    t.style = "Table Grid"
    t.alignment = WD_TABLE_ALIGNMENT.CENTER
    for i, h in enumerate(headers):
        cell = t.rows[0].cells[i]
        cell.text = ""
        p = cell.paragraphs[0]
        p.paragraph_format.space_before = Pt(1)
        p.paragraph_format.space_after = Pt(1)
        r = p.add_run(h)
        r.font.bold = True
        r.font.color.rgb = WColor(255, 255, 255)
        r.font.size = Pt(10.5)
        shade(cell, NAVY)
    for ri, row in enumerate(rows, start=1):
        cells = t.add_row().cells
        for ci, val in enumerate(row):
            cells[ci].text = ""
            p = cells[ci].paragraphs[0]
            p.paragraph_format.space_before = Pt(1)
            p.paragraph_format.space_after = Pt(1)
            r = p.add_run(str(val))
            r.font.size = Pt(10)
            if ci == 0:
                r.font.bold = True
            if zebra and ri % 2 == 0:
                shade(cells[ci], "F2F7FB")
    if widths:
        for row in t.rows:
            for i, w in enumerate(widths):
                row.cells[i].width = Inches(w)
    return t


def bullets_doc(doc, items, size=10.5):
    for b in items:
        p = doc.add_paragraph(style="List Bullet")
        p.paragraph_format.space_after = Pt(2)
        r = p.add_run(b)
        r.font.size = Pt(size)


def add_footer_pagenum(section, doc_title):
    sec = section
    sec.top_margin = Inches(0.8)
    sec.bottom_margin = Inches(0.8)
    sec.left_margin = Inches(0.9)
    sec.right_margin = Inches(0.9)
    sec.different_first_page_header_footer = True
    p = sec.footer.paragraphs[0]
    p.text = f"{PROJECT} — PRD · {doc_title}"
    for r in p.runs:
        r.font.size = Pt(8.5)
        r.font.color.rgb = WColor(0x99, 0xA7, 0xB4)
    from docx.enum.text import WD_TAB_ALIGNMENT

    p.paragraph_format.tab_stops.add_tab_stop(Inches(7.0), WD_TAB_ALIGNMENT.RIGHT)
    r = p.add_run("\t")
    fld = OxmlElement("w:fldSimple")
    fld.set(qn("w:instr"), "PAGE")
    run = OxmlElement("w:r")
    rPr = OxmlElement("w:rPr")
    sz = OxmlElement("w:sz")
    sz.set(qn("w:val"), "18")
    rPr.append(sz)
    t = OxmlElement("w:t")
    t.text = "1"
    run.append(rPr)
    run.append(t)
    fld.append(run)
    p._p.append(fld)


def cover_page(doc, title, subtitle, meta_rows):
    p = doc.add_paragraph()
    p.paragraph_format.space_after = Pt(2)
    r = p.add_run("HMSG (Hyundai Motorsport GmbH) — WEC Racing")
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.color.rgb = WColor(0, 170, 210)
    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(18)
    p.paragraph_format.space_after = Pt(6)
    r = p.add_run(title)
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = WColor(0, 44, 95)
    add_bottom_border(p, color=AMBER, size="18")
    p = doc.add_paragraph()
    r = p.add_run(subtitle)
    r.font.size = Pt(12.5)
    r.font.color.rgb = WColor(89, 89, 89)
    doc.add_paragraph().paragraph_format.space_after = Pt(2)
    t = doc.add_table(rows=0, cols=2)
    t.style = "Table Grid"
    for k, v in meta_rows:
        cells = t.add_row().cells
        cells[0].text = ""
        cells[1].text = ""
        r0 = cells[0].paragraphs[0].add_run(k)
        r0.font.bold = True
        r0.font.size = Pt(10)
        r0.font.color.rgb = WColor(255, 255, 255)
        shade(cells[0], NAVY)
        r1 = cells[1].paragraphs[0].add_run(v)
        r1.font.size = Pt(10)
        if k in ("Document", "Project"):
            r1.font.bold = True
        for c in t.rows[-1].cells:
            c.width = Inches(1.6 if c is cells[0] else 5.9)
    doc.add_page_break()


def build_word_doc(filename, title, subtitle, meta, sections):
    doc = Document()
    normal = doc.styles["Normal"]
    normal.font.name = "Calibri"
    normal.font.size = Pt(10.5)
    normal.font.color.rgb = WColor(0x33, 0x33, 0x33)
    for lvl in (1, 2, 3):
        style_heading(doc, lvl)
    cover_page(doc, title, subtitle, meta)
    add_footer_pagenum(doc.sections[0], title.split("—")[0].strip())
    for kind, *args in sections:
        if kind == "h1":
            doc.add_heading(args[0], level=1)
        elif kind == "h2":
            doc.add_heading(args[0], level=2)
        elif kind == "p":
            p = doc.add_paragraph(args[0])
            p.paragraph_format.space_after = Pt(6)
        elif kind == "bullets":
            bullets_doc(doc, args[0])
        elif kind == "table":
            doc_table(doc, *args)
        elif kind == "note":
            p = doc.add_paragraph()
            p.paragraph_format.space_before = Pt(4)
            r = p.add_run(args[0])
            r.font.size = Pt(10)
            r.font.italic = True
            r.font.color.rgb = (
                WColor(0x80, 0x00, 0x00)
                if args[1] == "warn"
                else WColor(0x59, 0x59, 0x59)
            )
    doc.save(BASE / filename)


META_BASE = [
    ("Document", ""),
    ("Project", "HMSG (Hyundai Motorsport GmbH) — WEC Racing"),
    ("Environment / Region", f"prd · {REGION}"),
    ("Stack", STACK),
    ("Version", "1.0"),
    ("Date", TODAY),
    ("Owner", "HAEE Cloud Team"),
    ("Status", "Draft — for review"),
]

# --- Administration Guide ---
meta = [["Document", "Administration Guide"]]
meta += META_BASE
meta += [
    ["Related", "sub Documents in the deliverables folder"],
    ["Baseline", "main @ def9a32"],
]
sections = [
    ("h1", "1. Purpose & scope"),
    (
        "p",
        "Day-2 guide for operating the HMSG WEC Racing production environment: the Windows "
        "MSSQL host, AWS Client VPN access, and (when enabled) the Site-to-Site VPN.",
    ),
    ("h1", "2. Access paths"),
    (
        "table",
        ["Path", "How to use", "Notes"],
        [
            [
                "AWS Client VPN",
                "Open the .ovpn file (AWS supplied client) and sign in with Entra ID SSO",
                "Primary access for developers (≤10 users)",
            ],
            [
                "RDP to EC2",
                "mstsc → MssqlPrivateIp (172.200.0.128–159 range) with sysadm",
                "Only reachable through Client VPN / on-prem — no public IP",
            ],
            [
                "SSM Session Manager",
                "EC2 console or aws ssm start-session; uses the instance profile",
                "Backup access when RDP is unavailable",
            ],
            [
                "AWS console / CLI",
                "AWS account 440027026402, eu-central-1, stack " + STACK,
                "Operations and troubleshooting",
            ],
        ],
        [1.6, 3.4, 2.5],
    ),
    ("h1", "3. Key resources"),
    (
        "table",
        ["Resource", "What it is"],
        [
            [
                "hmsg-rac-prd-vpc-ec1",
                "VPC 172.200.0.0/24 — dual-AZ public / private / DB subnets",
            ],
            [
                "hmsg-rac-prd-mssql01-ec2-ec1a",
                "Windows Server 2022 (m7i.2xlarge) — 250 GB root + 4 TB D: data volume",
            ],
            [
                "hmsg-rac-prd-clientvpn-ec1",
                "Client VPN endpoint — SAML (Entra ID) auth, UDP 443, split tunnel",
            ],
            [
                "hmsg-rac-prd-vpn-conn-ec1",
                "Site-to-Site VPN (optional) — static route to the on-prem network",
            ],
            [
                "hmsg-rac-prd-sysadm-password",
                "Secrets Manager secret — sysadm local administrator password",
            ],
        ],
        [3.1, 4.4],
    ),
    ("h1", "4. Operating rules"),
    (
        "bullets",
        [
            "UserData runs only at first boot: rotate sysadm by changing the Windows password on the host first, then update the secret.",
            "The DB subnet has no internet route by design — AWS APIs are only reachable through the VPC interface endpoints.",
            "Do not attach a public IP to the EC2 instance.",
            "Entra ID application and user management are client-owned after handover.",
            "Temporary implementation IAM permissions are removed at project completion.",
            "Client VPN → CloudWatch (90 days).",
        ],
    ),
    ("h1", "5. Routine checks"),
    (
        "table",
        ["Check", "Expected"],
        [
            [
                "Client VPN endpoint state",
                "available; clients connect via Entra ID SSO",
            ],
            [
                "4 TB data volume",
                "Mounted as D:, NTFS, label SQLData, at least 20% free for SQL growth",
            ],
            [
                "Port 1433 reachability",
                "Allowed from 10.200.0.0/22 (and on-prem CIDR when the S2S VPN exists)",
            ],
            [
                "Site-to-Site tunnel state",
                "UP on both ends; route advertised in pri/db route tables",
            ],
            ["Secrets Manager secret", "Matches the current Windows password"],
        ],
        [3.0, 4.5],
    ),
    ("h1", "6. Handover"),
    (
        "p",
        "At project completion: deliver the documentation set, remove temporary IAM access, "
        "transfer VPN user administration to HMGRS / the client, and obtain sign-off.",
    ),
    (
        "note",
        "Prepare for handover: export Security Group rules, route tables and endpoint IDs from the "
        "resource inventory workbook before transferring ownership.",
        "info",
    ),
]
build_word_doc(
    "Administration-Guide.docx",
    "Administration Guide",
    "Operating the HMSG WEC Racing production environment",
    meta,
    sections,
)

# --- Backup / Recovery / VPN ---
meta2 = [["Document", "Backup · MSSQL EC2 Recovery · VPN Configuration"]]
meta2 += META_BASE
sections2 = [
    ("h1", "1. Backup strategy"),
    (
        "p",
        "SQL Server is installed by the client; SQL backup scheduling, retention and media are therefore "
        "client-owned. This section defines the AWS-side duties and the coordination points.",
    ),
    (
        "table",
        ["Scope", "Who", "What to do"],
        [
            [
                "SQL backups",
                "Client (SQL admin)",
                "Regular full / differential / log backups per the agreed RPO; store on D: or an agreed share",
            ],
            [
                "EBS snapshots",
                "HMSG / ops",
                "Periodic snapshots of the root (250 GB) and data (4 TB) volumes — e.g. daily data-volume snapshot, weekly full",
            ],
            [
                "Client VPN logs",
                "AWS (automatic)",
                "CloudWatch log group hmsg-rac-prd-clientvpn-logs, 90-day retention",
            ],
        ],
        [1.5, 1.8, 4.2],
    ),
    (
        "bullets",
        [
            "Coordinate the snapshot window with the SQL backup window to keep the volumes consistent.",
            "Test restore at least once before go-live (see section 3).",
        ],
    ),
    ("h1", "2. MSSQL EC2 recovery"),
    ("h2", "2.1 Instance unhealthy / needs replacement"),
    (
        "bullets",
        [
            "Instance is protected (DeletionPolicy: Retain) — never delete it blindly; recover in place if the OS is intact.",
            "If replacement is required: restore the latest EBS snapshot to a new volume, or redeploy the CloudFormation stack and re-attach the data volume.",
            "Confirm the sysadm secret still matches; reset via EC2 Run Command / Session Manager if needed.",
        ],
    ),
    ("h2", "2.2 Volume failure (root or data)"),
    (
        "table",
        ["Case", "Recovery action"],
        [
            [
                "Root volume lost",
                "Redeploy the stack (or re-launch from the AMI), re-attach the retained data volume, ensure D: is mounted, verify the sysadm account.",
            ],
            [
                "Data volume lost",
                "Restore the latest data-volume snapshot as a new 4 TB volume, attach to D:, run DBCC CHECKDB before reopening access.",
            ],
            [
                "Both lost",
                "Full restore: new instance + restored data volume + SQL backup replay (full → diff → logs).",
            ],
        ],
        [2.0, 5.5],
    ),
    ("h2", "2.3 sysadm password recovery"),
    (
        "p",
        "Read the value from Secrets Manager, or decrypt the built-in Administrator password with the "
        "key pair (EC2 console → Get Windows password) and reset sysadm from there. Keep the secret "
        "in line with the account afterwards.",
    ),
    ("h1", "3. Recovery test"),
    (
        "bullets",
        [
            "Boot the restored instance and verify: D: mounted, SQL services start, 1433 reachable from a Client VPN client.",
            "Verify SSM Session Manager access (interface endpoints available).",
            "If the Site-to-Site VPN is deployed, verify the on-prem route in the pri/db route tables.",
            "Record the test result in the Project Completion Checklist evidence column.",
        ],
    ),
    ("h1", "4. VPN configuration"),
    ("h2", "4.1 AWS Client VPN"),
    (
        "table",
        ["Setting", "Value"],
        [
            ["Client CIDR", "10.200.0.0/22"],
            [
                "Authentication",
                "Federated SAML — Microsoft Entra ID (client-side enterprise app)",
            ],
            ["Transport / port", "UDP 443 — no TCP fallback"],
            ["Split tunnel", "Enabled"],
            [
                "Target networks",
                "pri-ec1a (172.200.0.64/27) · pri-ec1c (172.200.0.96/27)",
            ],
            ["Route", "172.200.0.0/24 via the pri-ec1a association"],
            [
                "Authorization",
                "All groups → whole VPC; restrict later via Entra group claims if required",
            ],
            ["Logging", "CloudWatch — 90 days"],
        ],
        [2.0, 5.5],
    ),
    ("h2", "4.2 SAML values for the client's Entra ID application"),
    (
        "table",
        ["Field", "Value"],
        [
            ["Identifier (Entity ID)", "urn:amazon:webservices:clientvpn"],
            [
                "Reply URL",
                "https://self-service.clientvpn.amazonaws.com/api/auth/sso/saml",
            ],
            [
                "Sign-on URL",
                "http://127.0.0.1:35001 (loopback listener on the user device)",
            ],
        ],
        [2.0, 5.5],
    ),
    ("h2", "4.3 Site-to-Site VPN (optional)"),
    (
        "bullets",
        [
            "Created only when CustomerGatewayIp and OnPremCidr are both supplied (add-only stack update).",
            "Customer gateway: ipsec.1, BGP ASN 65000. Amazon side ASN 64512. Static routing only.",
            "Download the tunnel configuration from the VPC console (Site-to-Site VPN connections) and hand it to the on-prem firewall team.",
            "If the race site moves (new public IP), update the stack — the VPN connection is replaced and must be re-configured on both sides.",
        ],
    ),
    (
        "note",
        "Never publish the tunnel pre-shared keys or the sysadm password in any shared document. "
        "Secrets stay in Secrets Manager and in the VPN consoles.",
        "warn",
    ),
]
build_word_doc(
    "Backup-MSSQL-EC2-Recovery-and-VPN-Configuration.docx",
    "Backup · MSSQL EC2 Recovery · VPN Configuration",
    "Recovery runbook and VPN reference for the HMSG WEC Racing production environment",
    meta2,
    sections2,
)

# ---------------------------------------------------------------------------
print("Deliverables regenerated in", BASE)
for p in sorted(BASE.iterdir()):
    if p.suffix in (".pptx", ".docx", ".xlsx"):
        print(f"  {p.name}: {p.stat().st_size} bytes")
